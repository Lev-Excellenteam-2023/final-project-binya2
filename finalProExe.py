import argparse
import asyncio
import json
import openai
from pptx import Presentation

def extract_text_from_slide(slide):
    return  " ".join([run.text for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]).strip()

async def sending_question_to_chat_gpt(presentation_text):
    # Complete a chat with GPT-3
    text = presentation_text + "can you summarize for me?"
    response = await asyncio.get_event_loop().run_in_executor(None, lambda: openai.ChatCompletion.create
    (model = "gpt-3.5-turbo",
     messages = [{"role": "user", "content": text}],
     timeout = 100000 # seconds
     ))
    return response.choices[0].message.content

def preparing_for_chat_question(file_path):
    pp_file = Presentation(file_path)
    actions = []
    for index, slide in enumerate(pp_file.slides, start=1):
        slide_text = extract_text_from_slide(slide)
        if slide_text:
            action = asyncio.create_task(sending_question_to_chat_gpt(slide_text))
            actions.append((index, slide_text, action))
    return actions

async def chat_gpt_answer(file_path):
    actions = preparing_for_chat_question(file_path)
    response_dict = {}
    for index, slide_text, action in actions:
        try:
            response = await action
        except Exception as e:
            response = f"Error occurred while processing slide {index}. Error message: {str(e)}"
        response_dict[f"response {index}"] = {"text": slide_text, "response": response}
    return json.dumps(response_dict, indent="\n")

async def main():
    user_path = argparse.ArgumentParser(description="Extract text from a PowerPoint file")
    user_path.add_argument("PowerPointFile", type=str, help="Path to the PowerPoint presentation file")
    file = user_path.parse_args()
    file_name = file.PowerPointFile.split("/")[-1].split(".")[0]
    with open(f"{file_name}.json", "w") as file_out:
        file_out.write(await chat_gpt_answer(file.PowerPointFile))
    print("You can read the result in the json file located in the folder of the presentation.")

if __name__ == "__main__":
    openai.api_key = "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"  # Replace with your API key
    asyncio.get_event_loop().run_until_complete(main())  # Run the main as async function