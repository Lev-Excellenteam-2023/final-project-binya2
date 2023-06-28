from pptx import Presentation
import argparse

def extract_text_from_slide(slide)->str:
    return  " ".join([run.text for shape in slide.shapes if shape.has_text_frame
                      for paragraph in shape.text_frame.paragraphs
                      for run in paragraph.runs]).strip()

def extext_text_from_pp_file(file_path)->list:
    pp_file = Presentation(file_path)
    slide_list = []
    for index, slide in enumerate(pp_file.slides, start=1):
        slide_text = extract_text_from_slide(slide)
        if slide_text:
            slide_list.append((index, slide_text))
    return slide_list


def extract_path_ftom_user()->str:
    user_path = argparse.ArgumentParser(description="Extract text from a PowerPoint file")
    user_path.add_argument("PowerPointFile", type=str, help="Path to the PowerPoint presentation file")
    file = user_path.parse_args()
    return file.PowerPointFile