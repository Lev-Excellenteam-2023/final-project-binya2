from pptx import Presentation
import argparse




def extract_text_from_slide(slide) -> str:
    """
    This module is used to extract text from a PowerPoint file.
    """
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape_text = extract_text_from_text_frame(shape.text_frame)
            text.append(shape_text)
    return " ".join(text).strip()

def extract_text_from_text_frame(text_frame) -> str:
    text = []
    for paragraph in text_frame.paragraphs:
        paragraph_text = extract_text_from_paragraph(paragraph)
        text.append(paragraph_text)
    return " ".join(text)

def extract_text_from_paragraph(paragraph) -> str:
    text = []
    for run in paragraph.runs:
        text.append(run.text)
    return " ".join(text)





def extext_text_from_pp_file(file_path) -> list:
    """
        This function gets a path to a PowerPoint file and returns a list of tuples.
    """
    pp_file = Presentation(file_path)
    slide_list = []
    for index, slide in enumerate(pp_file.slides, start=1):
        slide_text = extract_text_from_slide(slide)
        if slide_text:
            slide_list.append((index, slide_text))
    return slide_list





def extract_path_from_user() -> str:
    """
        This function gets a path to a PowerPoint file and returns a string.
    """
    user_path = argparse.ArgumentParser(description="Extract text from a PowerPoint file")
    user_path.add_argument("PowerPointFile", type=str, help="Path to the PowerPoint presentation file")
    file = user_path.parse_args()
    return file.PowerPointFile
